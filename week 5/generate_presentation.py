import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not found! Please run 'pip install python-pptx' and then run this script.")
    sys.exit(1)

# Color Palette Definitions
COLOR_BG_DARK = RGBColor(30, 41, 59)      # Slate 800 (Dark background)
COLOR_TEXT_LIGHT = RGBColor(248, 250, 252) # Slate 50 (White/Light text)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # Slate 50 (Light slide background)
COLOR_PRIMARY = RGBColor(30, 41, 59)      # Slate 800 (Dark text)
COLOR_TEAL = RGBColor(15, 118, 110)       # Teal 700 (Accent color)
COLOR_MUTED = RGBColor(100, 116, 139)     # Slate 500 (Subtitles, labels)
COLOR_SUCCESS = RGBColor(16, 185, 129)    # Emerald 500 (Conversion, positive metrics)
COLOR_CARD_BG = RGBColor(241, 245, 249)    # Slate 100 (Card background)

def add_title_slide(prs):
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Set dark background color using a full-screen rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()
    
    # Accent colored border line at the top
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = COLOR_TEAL
    top_line.line.fill.background()
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BANK MARKETING CAMPAIGN OPTIMIZATION"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Maximizing Deposit Subscriptions via Statistical Driver Analysis"
    p2.font.name = "Calibri"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEAL
    p2.font.italic = True
    p2.alignment = PP_ALIGN.LEFT
    
    # Subtitle Info (Program & Intern Details)
    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.2))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p3 = tf_info.paragraphs[0]
    p3.text = "AnalystLab Africa Data Analytics Internship Program"
    p3.font.name = "Calibri"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_TEXT_LIGHT
    
    p4 = tf_info.add_paragraph()
    p4.text = "Week 5 Case Study - Batch B (1st of June to 1st of August)\nPrepared by: Data Analyst Intern"
    p4.font.name = "Calibri"
    p4.font.size = Pt(12)
    p4.font.color.rgb = COLOR_MUTED

def add_header(slide, title_text, category_text="CASE STUDY ANALYSIS"):
    # Category tracker
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(8.8), Inches(0.3))
    tf_cat = cat_box.text_frame
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Calibri"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_TEAL
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.6))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY

def add_card(slide, left, top, width, height, title="", color=COLOR_CARD_BG):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = RGBColor(226, 232, 240) # Slate 200
    card.line.width = Pt(1)
    
    if title:
        tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Trebuchet MS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL
    
    return card

def add_bullet(tf, bold_prefix, text_content, pt_size=13, is_first=False):
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.space_after = Pt(8)
    
    run_bold = p.add_run()
    run_bold.text = bold_prefix + ": "
    run_bold.font.name = "Calibri"
    run_bold.font.size = Pt(pt_size)
    run_bold.font.bold = True
    run_bold.font.color.rgb = COLOR_PRIMARY
    
    run_text = p.add_run()
    run_text.text = text_content
    run_text.font.name = "Calibri"
    run_text.font.size = Pt(pt_size)
    run_text.font.color.rgb = COLOR_PRIMARY

def main():
    prs = Presentation()
    # Set to widescreen aspect ratio 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Business Problem
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide2, "The Business Problem Statement")
    # Left column - Context
    add_card(slide2, Inches(0.6), Inches(1.1), Inches(4.2), Inches(3.9), "Campaign Context")
    tb2_left = slide2.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(3.9), Inches(3.2))
    tf2_left = tb2_left.text_frame
    tf2_left.word_wrap = True
    add_bullet(tf2_left, "Core Objective", "Understand why clients subscribe to term deposits and identify high-conversion customer segments to boost ROI.", is_first=True)
    add_bullet(tf2_left, "Problem", "Direct telemarketing campaigns are expensive and suffer from low efficiency due to unsegmented cold-calling.", 13)
    add_bullet(tf2_left, "Impact", "Over-contacting customer pools causes campaign fatigue, brand irritation, and call-center resource drain.", 13)
    
    # Right column - Metrics Dashboard
    add_card(slide2, Inches(5.2), Inches(1.1), Inches(4.2), Inches(3.9), "Current Campaign Performance")
    
    # Big Number Callout
    stat_box = slide2.shapes.add_textbox(Inches(5.35), Inches(1.5), Inches(3.9), Inches(1.2))
    tf_stat = stat_box.text_frame
    p_stat = tf_stat.paragraphs[0]
    p_stat.text = "47.38%"
    p_stat.font.name = "Trebuchet MS"
    p_stat.font.size = Pt(48)
    p_stat.font.bold = True
    p_stat.font.color.rgb = COLOR_SUCCESS
    p_stat.alignment = PP_ALIGN.CENTER
    
    p_label = tf_stat.add_paragraph()
    p_label.text = "Average Subscription Rate (5,289/11,162 Subscribed)"
    p_label.font.name = "Calibri"
    p_label.font.size = Pt(11)
    p_label.font.bold = True
    p_label.font.color.rgb = COLOR_MUTED
    p_label.alignment = PP_ALIGN.CENTER
    
    # Key issues bullet points on right
    tb2_right = slide2.shapes.add_textbox(Inches(5.35), Inches(2.7), Inches(3.9), Inches(2.1))
    tf2_right = tb2_right.text_frame
    tf2_right.word_wrap = True
    add_bullet(tf2_right, "Volume Inefficiency", "May campaigns accounts for 25% of all calls, yet yield the lowest conversions (32.8%).", 12, is_first=True)
    add_bullet(tf2_right, "Unknown Methods", "Calling clients with unknown details drops conversion to only 22.6%.", 12)
    
    # Slide 3: Dataset & Methodology
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide3, "Dataset & Analytics Methodology")
    
    add_card(slide3, Inches(0.6), Inches(1.1), Inches(4.2), Inches(3.9), "Dataset Characteristics")
    tb3_left = slide3.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(3.9), Inches(3.2))
    tf3_left = tb3_left.text_frame
    tf3_left.word_wrap = True
    add_bullet(tf3_left, "Volume & Quality", "11,162 unique client outreach records with 100% complete data (no missing values).", 12, is_first=True)
    add_bullet(tf3_left, "Demographics", "Age, job profiles, marital status, and education levels.", 12)
    add_bullet(tf3_left, "Financial Profiles", "Account balances, housing/personal loan obligations, credit default status.", 12)
    add_bullet(tf3_left, "Campaign Detail", "Call duration, month, contact type, and previous marketing outcome status.", 12)
    
    add_card(slide3, Inches(5.2), Inches(1.1), Inches(4.2), Inches(3.9), "Methodological Workflow")
    tb3_right = slide3.shapes.add_textbox(Inches(5.35), Inches(1.6), Inches(3.9), Inches(3.2))
    tf3_right = tb3_right.text_frame
    tf3_right.word_wrap = True
    add_bullet(tf3_right, "Exploratory Data Analysis", "Calculated conditional probabilities to identify high-performing cohorts.", 12, is_first=True)
    add_bullet(tf3_right, "Feature Engineering", "Segmented account balance, age, and call count to discover non-linear trends.", 12)
    add_bullet(tf3_right, "Driver Significance", "Trained a Random Forest Classifier to identify features with the strongest predictive power.", 12)
    
    # Slide 4: Key Drivers
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide4, "Statistical Drivers of Subscription Success")
    
    # Add Feature Importance table
    table_shape = slide4.shapes.add_table(7, 3, Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(3.4)
    
    # Set headers
    headers = ["Predictive Feature", "Relative Importance", "Business Relevance"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TEAL
        for p in cell.text_frame.paragraphs:
            p.font.name = "Trebuchet MS"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_LIGHT
            
    # Table data
    data = [
        ("Call Duration (duration)", "36.87% (Rank 1)", "Strongest driver; length of pitch correlates to success"),
        ("Yearly Balance (balance)", "8.88% (Rank 2)", "Financial capacity determines savings capability"),
        ("Customer Age (age)", "8.40% (Rank 3)", "Non-linear U-shape (youth & retirees subscribe most)"),
        ("Outreach Month (month)", "8.32% (Rank 4)", "Strong seasonal trends, highlighting campaign fatigue"),
        ("Contact Channel (contact)", "4.48% (Rank 7)", "Cellular out-performs unknown methods (+32%)"),
        ("Prev Outcome (poutcome)", "3.34% (Rank 10)", "Previous successes represent high-conversion warm leads")
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(11)
                p.font.color.rgb = COLOR_PRIMARY
                if col_idx == 1:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_TEAL
                    
    # Slide 5: Duration & Campaign
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide5, "Operational Insights: Call Duration & Recontact Caps")
    
    add_card(slide5, Inches(0.6), Inches(1.1), Inches(4.2), Inches(3.9), "The Call Duration Effect")
    tb5_left = slide5.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(3.9), Inches(3.2))
    tf5_left = tb5_left.text_frame
    tf5_left.word_wrap = True
    add_bullet(tf5_left, "Immediate Rejection (<2 min)", "Only 9.02% of short calls convert. This represents cold-calling dismissals.", 12, is_first=True)
    add_bullet(tf5_left, "Interactive Pitch (3.5-5.5 min)", "A healthy subscription rate of 48.54% is achieved, showing customer engagement.", 12)
    add_bullet(tf5_left, "Product Deep-Dive (9.5+ min)", "84.70% conversion rate. Longer conversations allow detailed discussion of savings products.", 12)
    
    add_card(slide5, Inches(5.2), Inches(1.1), Inches(4.2), Inches(3.9), "Contact Frequency Diminishing Returns")
    tb5_right = slide5.shapes.add_textbox(Inches(5.35), Inches(1.6), Inches(3.9), Inches(3.2))
    tf5_right = tb5_right.text_frame
    tf5_right.word_wrap = True
    add_bullet(tf5_right, "First Attempt (1 call)", "Peak conversion at 53.38%. This is the prime opportunity window.", 12, is_first=True)
    add_bullet(tf5_right, "Second & Third Call", "Conversions drop slightly to 46.27% (2 calls) and 46.78% (3 calls).", 12)
    add_bullet(tf5_right, "Diminishing Returns (4-5 calls)", "Rates drop to 39.69%, indicating customer resistance.", 12)
    add_bullet(tf5_right, "Inefficient Over-Contact (10+ calls)", "Subscription rates crash to 22.38%. Wastes center resources and irritates clients.", 12)
    
    # Slide 6: Monthly Fatigue
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide6, "Temporal Insights: Monthly Fatigue vs. Off-Peak Gains")
    
    add_card(slide6, Inches(0.6), Inches(1.1), Inches(4.2), Inches(3.9), "The May Over-Saturation Bottleneck")
    tb6_left = slide6.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(3.9), Inches(3.2))
    tf6_left = tb6_left.text_frame
    tf6_left.word_wrap = True
    add_bullet(tf6_left, "Disproportionate Volume", "May represents 25.3% of total campaigns (2,824 calls).", 12, is_first=True)
    add_bullet(tf6_left, "Poorest Performance", "May yields the lowest subscription rate of just 32.76%.", 12)
    add_bullet(tf6_left, "Operational Drain", "High volumes of low-quality cold calls represent a massive sink of agent hours and acquisition costs.", 12)
    
    add_card(slide6, Inches(5.2), Inches(1.1), Inches(4.2), Inches(3.9), "The Off-Peak Opportunity")
    tb6_right = slide6.shapes.add_textbox(Inches(5.35), Inches(1.6), Inches(3.9), Inches(3.2))
    tf6_right = tb6_right.text_frame
    tf6_right.word_wrap = True
    add_bullet(tf6_right, "Transitional Peaks", "March (89.86%), September (84.33%), and October (82.40%) convert incredibly well.", 12, is_first=True)
    add_bullet(tf6_right, "Winter Surge", "December reaches a 90.91% conversion rate on low-volume, highly selective outreach.", 12)
    add_bullet(tf6_right, "Strategy", "Shift call center capacity away from May to these transitional months where customers are highly receptive.", 12)
    
    # Slide 7: Segmentation
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide7, "Customer Segments: Opportunities & Risks")
    
    add_card(slide7, Inches(0.6), Inches(1.1), Inches(4.2), Inches(3.9), "Target Profiles (Opportunity Segments)", RGBColor(240, 253, 250))
    tb7_left = slide7.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(3.9), Inches(3.2))
    tf7_left = tb7_left.text_frame
    tf7_left.word_wrap = True
    add_bullet(tf7_left, "Students & Retirees", "Students convert at 74.72% and Retirees convert at 66.32%. Highly receptive demographic.", 12, is_first=True)
    add_bullet(tf7_left, "High-Balance Clients", "Top balance quintile (>€1,708 balance) converts at 58.00%. Highly liquid capital.", 12)
    add_bullet(tf7_left, "Prior Successes", "Past campaign success converts at 91.32% in the current campaign. A primary source of warm leads.", 12)
    
    add_card(slide7, Inches(5.2), Inches(1.1), Inches(4.2), Inches(3.9), "Avoid Profiles (High Risk Segments)", RGBColor(254, 242, 242))
    tb7_right = slide7.shapes.add_textbox(Inches(5.35), Inches(1.6), Inches(3.9), Inches(3.2))
    tf7_right = tb7_right.text_frame
    tf7_right.word_wrap = True
    add_bullet(tf7_right, "Debtor Customers", "Clients with personal loans convert at 33.15% and housing loans at 36.64%. They have low disposable income.", 12, is_first=True)
    add_bullet(tf7_right, "Blue-Collar Workers", "Constitutes 1,944 calls but converts at only 36.42%. Heavy resource waste.", 12)
    add_bullet(tf7_right, "Unknown Contacts", "Outreach with unknown contact channels converts at only 22.59%.", 12)
    
    # Slide 8: Actionable Recommendations
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide8, "Strategic Recommendations for Campaign ROI", "EXECUTIVE RECOMMENDATIONS")
    
    # 3 major recommendation blocks
    # Rec 1
    add_card(slide8, Inches(0.6), Inches(1.1), Inches(8.8), Inches(1.1), "1. Reallocate Monthly Budgets and Re-balance Outreach Calendar")
    tb8_1 = slide8.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(8.5), Inches(0.7))
    tf8_1 = tb8_1.text_frame
    tf8_1.word_wrap = True
    p = tf8_1.paragraphs[0]
    p.text = "Shift agent resources out of over-saturated May campaigns. Re-route 50% of call volumes to March, September, October, and December where conversion is high (82%-90%)."
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_PRIMARY
    
    # Rec 2
    add_card(slide8, Inches(0.6), Inches(2.35), Inches(8.8), Inches(1.1), "2. Implement Hard Call Capping (Maximum 3 Attempts)")
    tb8_2 = slide8.shapes.add_textbox(Inches(0.75), Inches(2.65), Inches(8.5), Inches(0.7))
    tf8_2 = tb8_2.text_frame
    tf8_2.word_wrap = True
    p = tf8_2.paragraphs[0]
    p.text = "Cap attempts at 3. Re-contacting customers beyond 3 attempts yields sharp diminishing returns (<31%) and damages customer goodwill. Shift idle hours to warm leads."
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_PRIMARY
    
    # Rec 3
    add_card(slide8, Inches(0.6), Inches(3.6), Inches(8.8), Inches(1.1), "3. Launch Segmented Offer Pushing and Channel Modernization")
    tb8_3 = slide8.shapes.add_textbox(Inches(0.75), Inches(3.9), Inches(8.5), Inches(0.7))
    tf8_3 = tb8_3.text_frame
    tf8_3.word_wrap = True
    p = tf8_3.paragraphs[0]
    p.text = "Prioritize Students/Retirees with tailored terms (low risk, predictable return). Restructure blue-collar pitches to focus on cash flow. Enrich customer databases to eliminate unknown contact channels."
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_PRIMARY
    
    # Slide 9: Conclusion
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark slide background for conclusion
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = COLOR_BG_DARK
    bg9.line.fill.background()
    
    # Title
    title_box9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.8))
    p_title9 = title_box9.text_frame.paragraphs[0]
    p_title9.text = "Key Execution Takeaways"
    p_title9.font.name = "Trebuchet MS"
    p_title9.font.size = Pt(28)
    p_title9.font.bold = True
    p_title9.font.color.rgb = COLOR_TEXT_LIGHT
    
    # Core conclusions
    tb9_body = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.0))
    tf9_body = tb9_body.text_frame
    tf9_body.word_wrap = True
    
    p = tf9_body.paragraphs[0]
    p.space_after = Pt(12)
    r = p.add_run()
    r.text = "Quality Over Volume: "
    r.font.bold = True
    r.font.color.rgb = COLOR_TEAL
    r.font.size = Pt(14)
    r2 = p.add_run()
    r2.text = "Shifting from unsegmented high-volume cold-calling to targeted demographics can yield the same deposit volumes with up to 40% lower contact volume."
    r2.font.color.rgb = COLOR_TEXT_LIGHT
    r2.font.size = Pt(14)
    
    p2 = tf9_body.add_paragraph()
    p2.space_after = Pt(12)
    r = p2.add_run()
    r.text = "Enhance Agent Performance: "
    r.font.bold = True
    r.font.color.rgb = COLOR_TEAL
    r.font.size = Pt(14)
    r2 = p2.add_run()
    r2.text = "Training agents to sustain conversational engagement past the critical 3-minute mark yields an immediate 5x increase in conversion rate."
    r2.font.color.rgb = COLOR_TEXT_LIGHT
    r2.font.size = Pt(14)
    
    p3 = tf9_body.add_paragraph()
    r = p3.add_run()
    r.text = "Prioritize High Value Leads: "
    r.font.bold = True
    r.font.color.rgb = COLOR_TEAL
    r.font.size = Pt(14)
    r2 = p3.add_run()
    r2.text = "Prioritize clients with a history of success (91.3% conversion rate) and liquid wealth (>€1,708 balance) to optimize direct marketing efforts."
    r2.font.color.rgb = COLOR_TEXT_LIGHT
    r2.font.size = Pt(14)

    prs.save("bank_marketing_case_study.pptx")
    print("PowerPoint presentation 'bank_marketing_case_study.pptx' generated successfully!")

if __name__ == "__main__":
    main()
